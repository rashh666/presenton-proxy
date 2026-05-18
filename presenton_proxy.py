#!/usr/bin/env python3
"""
presenton_proxy.py — Production-grade multi-agent proxy for Presenton (v6.0.0)
"Native Sequential" architecture: llama-server processes run directly on the host,
one per AMD RX 9060 XT GPU (ROCm), with a VRAM-clearing idle watchdog.
"""

from __future__ import annotations

import asyncio
import base64
import json
import logging
import re
import time
import uuid
import os
from contextlib import asynccontextmanager
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

import httpx
from fastapi import Depends, FastAPI, Header, HTTPException, Query, Request, Security
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, Response
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from fastapi.staticfiles import StaticFiles
from prometheus_client import CONTENT_TYPE_LATEST, Counter, Histogram, generate_latest
from pydantic import BaseModel, Field, ValidationError
from pydantic_settings import BaseSettings
from tenacity import RetryError, retry, retry_if_exception, stop_after_attempt, wait_exponential
import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor
from svg_generator import generate_svg_pictograph

# ---------------------------------------------------------------------------
# Optional faster JSON
try:
    import orjson as _fj

    def json_dumps(obj: Any) -> str:
        return _fj.dumps(obj).decode()

    def json_loads(s: Union[str, bytes]) -> Any:
        return _fj.loads(s)

except ImportError:
    import json as _sj

    def json_dumps(obj: Any) -> str:
        return _sj.dumps(obj, ensure_ascii=False)

    def json_loads(s: Union[str, bytes]) -> Any:
        return _sj.loads(s)

# Optional schema validation
try:
    import jsonschema
    JSONSCHEMA_AVAILABLE = True
except ImportError:
    JSONSCHEMA_AVAILABLE = False

# ---------------------------------------------------------------------------
# Settings
class Settings(BaseSettings):
    # Role-isolated endpoints
    reasoner_url: str = "http://localhost:8081/v1/chat/completions"
    reasoner_health_url: str = "http://localhost:8081/health"
    coder_url: str = "http://localhost:8082/v1/chat/completions"
    coder_health_url: str = "http://localhost:8082/health"

    # Model selection
    reasoner_model_key: str = "gemma3"      # options: gemma3, dagger
    coder_model_key: str = "codegeex4"      # options: codegeex4, gemma2

    # Hallucination mode
    hallucination_mode: str = "tight"
    temp_tight: float = 0.0
    top_p_tight: float = 0.1
    temp_accommodative: float = 0.8
    top_p_accommodative: float = 0.95

    # Default persona & narrative
    default_persona: str = "balanced"
    default_narrative: str = "problem-solution"

    # Request limits
    max_schema_size: int = 100_000
    max_messages: int = 50
    max_body_bytes: int = 1_000_000

    # Image generation (optional)
    image_gen_enabled: bool = False
    image_gen_url: str = "http://localhost:7860/sdapi/v1/txt2img"
    image_gen_width: int = 1024
    image_gen_height: int = 576
    image_gen_steps: int = 4
    image_gen_cfg_scale: float = 1.5
    image_gen_negative: str = "cartoon, drawing, anime, blurry, low quality, distorted text, messy"
    image_placeholder_url: str = "/images/placeholder.png"
    image_output_dir: Path = Path("./generated_images")
    image_gen_prefix: str = "professional corporate photography"
    image_prompt_suffix: str = "cinematic lighting, 8k, bokeh, sharp focus"
    image_gen_timeout: int = 60
    image_gen_semaphore: int = 2

    # Reflection / critic loop
    reflection_enabled: bool = False
    reflection_max_iterations: int = 1

    # Visual co-pilot temperature
    visual_copilot_temp: float = 0.6
    visual_copilot_top_p: float = 0.9

    # CORS (comma-separated)
    cors_origins: str = ""

    # API authentication (optional)
    api_key: Optional[str] = None

    # Native process management
    llama_server_bin: str = "/home/rashid/llama.cpp/build/bin/llama-server"
    models_host_path: str = os.environ.get("MIA_MODELS_PATH", "/models")
    process_startup_timeout: int = 120   # seconds to wait for health after spawn
    idle_timeout: int = 300              # seconds of inactivity before VRAM flush
    llm_timeout: int = 600               # httpx read timeout for LLM inference calls

    # Unified single-model mode: one llama-server spanning both GPUs on port 8081.
    # Both reasoner and coder phases hit the same endpoint.
    unified_single_model: bool = True

    # PPTX output directory
    pptx_output_dir: Path = Path("./presentations")

    # Image job TTL (seconds)
    image_job_ttl: int = 3600

    # Keys to search for slides list (top-level, then depth-1)
    slides_list_keys: List[str] = ["slides", "content", "sections", "outline"]

    # Skip these keys when searching depth-1
    skip_keys_depth1: List[str] = ["_meta", "_visual_hints", "presenton"]

    class Config:
        env_file = ".env"
        env_file_encoding = "utf-8"


settings = Settings()

# ---------------------------------------------------------------------------
# Logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
)
_base_logger = logging.getLogger("presenton-proxy")


class RequestLoggerAdapter(logging.LoggerAdapter):
    def process(self, msg: str, kwargs: dict) -> Tuple[str, dict]:
        return "[%s] %s" % (self.extra.get("request_id", "no-id"), msg), kwargs


def make_logger(request_id: str) -> RequestLoggerAdapter:
    return RequestLoggerAdapter(_base_logger, {"request_id": request_id})


# ---------------------------------------------------------------------------
# Metrics (prometheus_client)
_REQUESTS_TOTAL = Counter(
    "presenton_requests_total",
    "Total proxy requests",
    ["endpoint", "status"],
)
_REQUEST_LATENCY = Histogram(
    "presenton_request_duration_seconds",
    "End-to-end request latency in seconds",
    buckets=[0.1, 0.5, 1.0, 2.0, 5.0, 10.0, 30.0, 60.0, 120.0],
)

# ---------------------------------------------------------------------------
# Enums & constants

class HallucinationMode(str, Enum):
    TIGHT = "tight"
    ACCOMMODATIVE = "accommodative"


class PipelineMode(str, Enum):
    FULL = "full"
    OUTLINE_ONLY = "outline-only"


class NarrativeStructure(str, Enum):
    PROBLEM_SOLUTION = "problem-solution"
    HERO_JOURNEY = "hero-journey"
    INVESTOR_PITCH = "investor-pitch"
    EXECUTIVE_BRIEFING = "executive-briefing"
    CASE_STUDY = "case-study"


NARRATIVE_INSTRUCTIONS: Dict[NarrativeStructure, str] = {
    NarrativeStructure.PROBLEM_SOLUTION: (
        "Structure: Hook -> Problem -> Consequences -> Solution -> Proof -> CTA. "
        "Story beats per slide must be one of: hook, tension, problem, consequence, "
        "solution, evidence, proof, CTA."
    ),
    NarrativeStructure.HERO_JOURNEY: (
        "Structure: Ordinary World -> Call -> Challenge -> Transformation -> Return. "
        "Story beats: ordinary, call, refusal, mentor, challenge, revelation, transformation, return."
    ),
    NarrativeStructure.INVESTOR_PITCH: (
        "Structure: Vision -> Problem -> Market -> Solution -> Traction -> Team -> Ask. "
        "Story beats: vision, problem, market, solution, traction, team, financials, ask."
    ),
    NarrativeStructure.EXECUTIVE_BRIEFING: (
        "Structure: Situation -> Complication -> Resolution -> Recommendation. "
        "Story beats: situation, complication, implication, resolution, recommendation."
    ),
    NarrativeStructure.CASE_STUDY: (
        "Structure: Client Context -> Challenge -> Approach -> Results -> Learnings. "
        "Story beats: context, challenge, approach, execution, results, learnings."
    ),
}

# ---------------------------------------------------------------------------
# Persona profiles
PERSONA_PROFILES: Dict[str, Dict[str, Any]] = {
    "balanced": {
        "tone": "professional and approachable",
        "sentence_style": "clear, medium-length sentences with active voice",
        "slide_density": "moderate - 4-6 bullets, concise",
        "image_style": "professional corporate photography",
        "layout_bias": "bullet-list",
        "visual_palette": ["#1E3A5F", "#2D9CDB", "#F2F2F2"],
        "rhetorical_devices": ["parallelism", "rule of three"],
        "typography": {"title": "Inter", "body": "Source Sans Pro"},
        "reflection_gate": False,
    },
    "steve_jobs": {
        "tone": "visionary, minimalist, emotionally resonant",
        "sentence_style": "short punchy statements. One idea per slide. Dramatic pauses.",
        "slide_density": "sparse - 1-3 lines max. Let silence speak.",
        "image_style": "cinematic product photography, black background, single hero object",
        "layout_bias": "full-image",
        "visual_palette": ["#000000", "#FFFFFF", "#C0C0C0"],
        "rhetorical_devices": ["rule of three", "anaphora", "contrast", "the pause"],
        "typography": {"title": "Helvetica Neue", "body": "Helvetica Neue Light"},
        "reflection_gate": True,
    },
    "ted_talk": {
        "tone": "warm, story-driven, intellectually curious",
        "sentence_style": "conversational, first-person anecdotes, rhetorical questions",
        "slide_density": "light - slides support the speaker, not replace them",
        "image_style": "candid human moments, natural lighting, diverse subjects",
        "layout_bias": "full-image",
        "visual_palette": ["#EB0028", "#FFFFFF", "#1A1A1A"],
        "rhetorical_devices": ["anecdote", "rhetorical question", "callback", "contrast"],
        "typography": {"title": "Georgia", "body": "Georgia"},
        "reflection_gate": True,
    },
    "investor_pitch": {
        "tone": "confident, data-driven, urgency without desperation",
        "sentence_style": "declarative statements, numbers front-loaded, no hedging",
        "slide_density": "tight - one key metric or claim per slide",
        "image_style": "clean infographics, charts, market maps",
        "layout_bias": "stat-highlight",
        "visual_palette": ["#0F172A", "#38BDF8", "#22D3EE"],
        "rhetorical_devices": ["statistics", "social proof", "future pacing"],
        "typography": {"title": "Inter", "body": "Inter"},
        "reflection_gate": True,
    },
    "academic_lecture": {
        "tone": "precise, measured, citation-aware",
        "sentence_style": "passive voice acceptable; definitions before claims; hedge appropriately",
        "slide_density": "dense - up to 8 bullets; methodology detail matters",
        "image_style": "diagrams, flowcharts, data plots, microscopy",
        "layout_bias": "two-column",
        "visual_palette": ["#1B2A3B", "#4A90D9", "#F5F5F5"],
        "rhetorical_devices": ["definition", "citation", "comparative analysis"],
        "typography": {"title": "Palatino", "body": "Times New Roman"},
        "reflection_gate": False,
    },
    "cyberpunk_founder": {
        "tone": "bold, irreverent, future-obsessed",
        "sentence_style": "fragments welcome. Tech slang. Provocative openers.",
        "slide_density": "sparse - punchy, maximum contrast",
        "image_style": "neon-lit cyberpunk cityscape, glitch art, holographic UI",
        "layout_bias": "full-image",
        "visual_palette": ["#0D0D0D", "#00FF9C", "#FF00FF"],
        "rhetorical_devices": ["provocation", "contrast", "future pacing"],
        "typography": {"title": "Space Grotesk", "body": "JetBrains Mono"},
        "reflection_gate": False,
    },
}

# ---------------------------------------------------------------------------
# Layout intelligence
class LayoutType(str, Enum):
    TITLE_SLIDE    = "title-slide"
    TWO_COLUMN     = "two-column"
    BULLET_LIST    = "bullet-list"
    FULL_IMAGE     = "full-image"
    QUOTE          = "quote"
    STAT_HIGHLIGHT = "stat-highlight"
    TIMELINE       = "timeline"
    COMPARISON     = "comparison"


LAYOUT_DEFINITIONS: Dict[str, Dict[str, Any]] = {
    LayoutType.TITLE_SLIDE.value:    {"image_weight": 0.0,  "text_weight": 1.0,  "max_bullets": 0},
    LayoutType.TWO_COLUMN.value:     {"image_weight": 0.45, "text_weight": 0.55, "max_bullets": 5},
    LayoutType.BULLET_LIST.value:    {"image_weight": 0.0,  "text_weight": 1.0,  "max_bullets": 7},
    LayoutType.FULL_IMAGE.value:     {"image_weight": 0.85, "text_weight": 0.15, "max_bullets": 2},
    LayoutType.QUOTE.value:          {"image_weight": 0.0,  "text_weight": 1.0,  "max_bullets": 1},
    LayoutType.STAT_HIGHLIGHT.value: {"image_weight": 0.2,  "text_weight": 0.8,  "max_bullets": 3},
    LayoutType.TIMELINE.value:       {"image_weight": 0.1,  "text_weight": 0.9,  "max_bullets": 6},
    LayoutType.COMPARISON.value:     {"image_weight": 0.1,  "text_weight": 0.9,  "max_bullets": 6},
}


class Layout(BaseModel):
    type: LayoutType = LayoutType.BULLET_LIST
    image_weight: float = Field(default=0.0, ge=0.0, le=1.0)
    text_weight: float = Field(default=1.0, ge=0.0, le=1.0)

    @classmethod
    def from_type(cls, layout_type: LayoutType) -> "Layout":
        defn = LAYOUT_DEFINITIONS[layout_type.value]
        return cls(
            type=layout_type,
            image_weight=defn["image_weight"],
            text_weight=defn["text_weight"],
        )


class VisualHints(BaseModel):
    icon: Optional[str] = None
    layout: Layout = Field(default_factory=Layout)
    image_prompt: Optional[str] = None
    enhanced_prompt: Optional[str] = None
    style_tags: List[str] = Field(default_factory=list)
    theme: Optional[str] = None
    font_pair: Optional[Dict[str, str]] = None
    colors: List[str] = Field(default_factory=list)


# ---------------------------------------------------------------------------
# Pydantic request/response models
class ChatMessage(BaseModel):
    role: str
    content: str


class ChatCompletionRequest(BaseModel):
    messages: List[ChatMessage]
    response_format: Dict[str, Any] = Field(default_factory=dict)


class ExpandRequest(BaseModel):
    outline: List[Dict[str, Any]]
    response_format: Dict[str, Any] = Field(default_factory=dict)
    persona: str = "balanced"
    narrative_structure: str = "problem-solution"


# ---------------------------------------------------------------------------
# Error taxonomy
class ErrorCode(str, Enum):
    SCHEMA_VALIDATION_FAILED = "SCHEMA_VALIDATION_FAILED"
    SUPERVISOR_FAILED        = "SUPERVISOR_FAILED"
    WORKER_FAILED            = "WORKER_FAILED"
    CRITIC_FAILED            = "CRITIC_FAILED"
    IMAGE_JOB_NOT_FOUND      = "IMAGE_JOB_NOT_FOUND"
    REQUEST_TOO_LARGE        = "REQUEST_TOO_LARGE"
    INVALID_REQUEST          = "INVALID_REQUEST"
    UNSUPPORTED_FORMAT       = "UNSUPPORTED_FORMAT"
    INTERNAL_ERROR           = "INTERNAL_ERROR"
    UNAUTHORIZED             = "UNAUTHORIZED"
    NOT_FOUND                = "NOT_FOUND"


def error_response(
    status_code: int,
    code: ErrorCode,
    message: str,
    request_id: str,
) -> JSONResponse:
    return JSONResponse(
        status_code=status_code,
        content={"error": {"code": code.value, "message": message, "request_id": request_id}},
        headers={"X-Request-ID": request_id},
    )


def raise_api_error(code: ErrorCode, message: str, http_status: int) -> None:
    raise HTTPException(status_code=http_status, detail={"code": code.value, "message": message})


def _http_exc_to_response(exc: HTTPException, request_id: str) -> JSONResponse:
    detail = exc.detail
    if isinstance(detail, dict):
        try:
            code = ErrorCode(detail.get("code", ErrorCode.INTERNAL_ERROR.value))
        except ValueError:
            code = ErrorCode.INTERNAL_ERROR
        msg = detail.get("message", "Unknown error")
    else:
        code = ErrorCode.INTERNAL_ERROR
        msg = str(detail)
    return error_response(exc.status_code, code, msg, request_id)


# ---------------------------------------------------------------------------
# JSON utilities
def robust_json_loads(text: str) -> Dict[str, Any]:
    """Extract JSON from model output, handling markdown fences and trailing text."""
    text = text.strip()
    text = re.sub(r"```(?:json)?\s*\n?", "", text).strip()

    decoder = json.JSONDecoder()
    # Strategy 1: raw_decode from start
    try:
        obj, _ = decoder.raw_decode(text)
        if isinstance(obj, dict):
            return obj
    except (json.JSONDecodeError, ValueError):
        pass

    # Strategy 2: find first '{' and raw_decode from there
    start = text.find("{")
    if start != -1:
        try:
            obj, _ = decoder.raw_decode(text[start:])
            if isinstance(obj, dict):
                return obj
        except (json.JSONDecodeError, ValueError):
            pass

    # Strategy 3: brace-counting fallback
    if start != -1:
        brace_count = 0
        for i, ch in enumerate(text[start:], start=start):
            if ch == "{":
                brace_count += 1
            elif ch == "}":
                brace_count -= 1
                if brace_count == 0:
                    candidate = text[start:i+1]
                    try:
                        obj = json_loads(candidate)
                        if isinstance(obj, dict):
                            return obj
                    except (json.JSONDecodeError, ValueError):
                        break
    raise ValueError("Could not extract a valid JSON object from model output.")


def validate_against_schema(data: Dict, schema: Dict) -> None:
    if not JSONSCHEMA_AVAILABLE:
        return
    validator = jsonschema.Draft7Validator(schema)
    errors = sorted(validator.iter_errors(data), key=lambda e: list(e.path))
    if errors:
        msg = "\n".join(f"- {e.message} (path: {list(e.path)})" for e in errors[:5])
        raise ValueError(f"Schema validation failed:\n{msg}")


# ---------------------------------------------------------------------------
# Prompt sanitisation
_UNSAFE_PATTERNS = [
    re.compile(r"ignore\s+(previous|all|above)\s+instructions", re.IGNORECASE),
    re.compile(r"disregard\s+system\s+prompt", re.IGNORECASE),
    re.compile(r"forget\s+(your|all)\s+(previous\s+)?instructions", re.IGNORECASE),
    re.compile(r"output\s+only\s+json", re.IGNORECASE),
]


def sanitize_prompt(text: str) -> str:
    for pat in _UNSAFE_PATTERNS:
        text = pat.sub("[blocked]", text)
    return text


def _is_transient(exc: Exception) -> bool:
    if isinstance(exc, httpx.TimeoutException):
        return True
    if isinstance(exc, httpx.HTTPStatusError):
        return exc.response.status_code in (429, 502, 503, 504)
    return False


def parse_enum_ci(enum_cls: type, value: str, default: Any, log_warning: bool = True) -> Any:
    """Case-insensitive enum lookup with logged fallback."""
    try:
        return enum_cls(value.lower())
    except ValueError:
        if log_warning:
            _base_logger.warning(
                "Unknown value '%s' for %s – falling back to '%s'",
                value, enum_cls.__name__, default
            )
        return default


_http_bearer = HTTPBearer(auto_error=False)


async def require_api_key(
    credentials: Optional[HTTPAuthorizationCredentials] = Security(_http_bearer),
) -> None:
    if settings.api_key is None:
        return
    token = credentials.credentials if credentials else None
    if not token or token != settings.api_key:
        raise HTTPException(
            status_code=401,
            detail={"code": ErrorCode.UNAUTHORIZED.value, "message": "Invalid or missing API key."},
        )


# ---------------------------------------------------------------------------
# Circuit Breaker
class CircuitBreakerOpen(Exception):
    pass


class CircuitBreaker:
    """Three-state circuit breaker (CLOSED → OPEN → HALF_OPEN → CLOSED).

    Trips to OPEN after `failure_threshold` consecutive failures; re-tests
    after `reset_timeout` seconds via a single HALF_OPEN probe.
    """

    CLOSED = "closed"
    OPEN = "open"
    HALF_OPEN = "half_open"

    def __init__(self, failure_threshold: int = 5, reset_timeout: float = 30.0) -> None:
        self._state = self.CLOSED
        self._failure_count = 0
        self._failure_threshold = failure_threshold
        self._reset_timeout = reset_timeout
        self._last_failure_time: Optional[float] = None
        self._lock = asyncio.Lock()

    @property
    def state(self) -> str:
        return self._state

    async def call(self, func, *args, **kwargs):
        async with self._lock:
            if self._state == self.OPEN:
                elapsed = (
                    time.time() - self._last_failure_time
                    if self._last_failure_time is not None
                    else 0.0
                )
                if elapsed >= self._reset_timeout:
                    self._state = self.HALF_OPEN
                else:
                    raise CircuitBreakerOpen(
                        f"Circuit open – downstream stalled. "
                        f"Retry in {self._reset_timeout - elapsed:.0f}s."
                    )
        try:
            result = await func(*args, **kwargs)
            async with self._lock:
                self._failure_count = 0
                self._state = self.CLOSED
            return result
        except CircuitBreakerOpen:
            raise
        except Exception:
            async with self._lock:
                self._failure_count += 1
                self._last_failure_time = time.time()
                if self._failure_count >= self._failure_threshold:
                    if self._state != self.OPEN:
                        _base_logger.warning(
                            "Circuit OPEN after %d consecutive failures on %s",
                            self._failure_count, self,
                        )
                    self._state = self.OPEN
            raise


# ---------------------------------------------------------------------------
# LLM HTTP client
class LLMClient:
    def __init__(
        self,
        base_url: str,
        timeout: int = 120,
        circuit: Optional[CircuitBreaker] = None,
    ) -> None:
        self.base_url = base_url
        self._client = httpx.AsyncClient(
            timeout=timeout,
            limits=httpx.Limits(max_keepalive_connections=10, max_connections=20),
        )
        self._circuit = circuit or CircuitBreaker()

    @retry(
        retry=retry_if_exception(_is_transient),
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=1, max=10),
        reraise=True,
    )
    async def generate(
        self,
        messages: List[Dict],
        temperature: float,
        top_p: float,
        response_format: Dict,
    ) -> Dict:
        payload = {
            "messages": messages,
            "temperature": temperature,
            "top_p": top_p,
            "response_format": response_format,
        }
        return await self._circuit.call(self._post, payload)

    async def _post(self, payload: Dict) -> Dict:
        resp = await self._client.post(self.base_url, json=payload)
        resp.raise_for_status()
        return resp.json()

    async def health(self, health_url: str) -> bool:
        try:
            async with httpx.AsyncClient(timeout=3) as c:
                r = await c.get(health_url)
                return r.status_code == 200
        except Exception:
            return False

    async def close(self) -> None:
        await self._client.aclose()


# ---------------------------------------------------------------------------
# Native Model Manager
class NativeModelManager:
    """Manages two llama-server host processes: reasoner (GPU 0 / port 8081)
    and coder (GPU 1 / port 8082).  An idle watchdog terminates either process
    after settings.idle_timeout seconds of inactivity to flush VRAM; the next
    request triggers an on-demand restart.
    """

    _ROLE_CONFIG: Dict[str, Dict[str, Any]] = {
        "reasoner": {"port": 8081, "gpu_idx": 0, "health_url": "http://localhost:8081/health"},
        "coder":    {"port": 8082, "gpu_idx": 1, "health_url": "http://localhost:8082/health"},
    }

    def __init__(self) -> None:
        self._procs: Dict[str, Optional[asyncio.subprocess.Process]] = {
            "reasoner": None, "coder": None,
        }
        self._locks: Dict[str, asyncio.Lock] = {
            "reasoner": asyncio.Lock(), "coder": asyncio.Lock(),
        }
        self._last_activity: Dict[str, float] = {
            "reasoner": time.time(), "coder": time.time(),
        }
        self._watchdog_task: Optional[asyncio.Task] = None

    # ------------------------------------------------------------------
    # Public interface

    async def start_all(self) -> bool:
        if settings.unified_single_model:
            _base_logger.info(
                "Native model manager: unified mode — single llama-server on port 8081 "
                "spanning both GPUs (HIP_VISIBLE_DEVICES=0,1)."
            )
            ok = await self._start_role("reasoner")
            if ok:
                _base_logger.info("Unified model process is healthy.")
            else:
                _base_logger.warning("Unified model process failed to become healthy.")
            self._watchdog_task = asyncio.create_task(self._watchdog())
            return ok
        else:
            _base_logger.info("Native model manager: spawning both llama-server processes...")
            reasoner_ok = await self._start_role("reasoner")
            coder_ok = await self._start_role("coder")
            if reasoner_ok and coder_ok:
                _base_logger.info("Both model processes are healthy.")
            else:
                _base_logger.warning("One or both model processes failed to become healthy.")
            self._watchdog_task = asyncio.create_task(self._watchdog())
            return reasoner_ok and coder_ok

    async def ensure_role_ready(self, role: str) -> bool:
        """Touch activity timer; restart the process if it has exited or is unresponsive.
        In unified mode, both roles share the single 'reasoner' process."""
        # In unified mode everything routes through "reasoner"
        effective_role = "reasoner" if settings.unified_single_model else role
        self._last_activity[effective_role] = time.time()
        config = self._ROLE_CONFIG[effective_role]
        proc = self._procs.get(effective_role)

        if proc is not None and proc.returncode is None:
            if await self._wait_for_health(config["health_url"], timeout=5):
                return True

        _base_logger.warning("%s is unresponsive or exited — restarting on demand", effective_role)
        return await self._start_role(effective_role)

    def process_status(self) -> Dict[str, str]:
        if settings.unified_single_model:
            proc = self._procs.get("reasoner")
            status = (
                "idle (terminated)" if proc is None
                else "running" if proc.returncode is None
                else f"exited({proc.returncode})"
            )
            return {"unified (GPU 0+1 / port 8081)": status}
        result = {}
        for role, proc in self._procs.items():
            result[role] = (
                "idle (terminated)" if proc is None
                else "running" if proc.returncode is None
                else f"exited({proc.returncode})"
            )
        return result

    def touch(self, role: str) -> None:
        """Heartbeat: reset the idle timer for role so the watchdog won't evict it.
        In unified mode always updates the 'reasoner' timestamp regardless of role."""
        effective = "reasoner" if settings.unified_single_model else role
        self._last_activity[effective] = time.time()

    async def shutdown(self) -> None:
        if self._watchdog_task:
            self._watchdog_task.cancel()
            try:
                await self._watchdog_task
            except asyncio.CancelledError:
                pass
        for role in ("reasoner", "coder"):
            await self._terminate_role(role)
        _base_logger.info("All llama-server processes terminated.")

    # ------------------------------------------------------------------
    # Internal helpers

    def _build_cmd(self, role: str) -> List[str]:
        config = self._ROLE_CONFIG[role]
        model_key = (
            settings.reasoner_model_key if role == "reasoner" else settings.coder_model_key
        )
        base = settings.models_host_path
        model_map: Dict[str, Tuple[str, str]] = {
            "gemma3":    (f"{base}/reasoner/gemma3.gguf",  "8192"),
            "dagger":    (f"{base}/reasoner/dagger.gguf",  "8192"),
            "codegeex4": (f"{base}/coder/codegeex4.gguf",  "16384"),
            "gemma2":    (f"{base}/coder/verifier.gguf",   "8192"),
        }
        model_path, ctx_size = model_map.get(model_key, (f"{base}/{model_key}.gguf", "8192"))
        return [
            settings.llama_server_bin,
            "--host", "127.0.0.1",
            "--port", str(config["port"]),
            "-m", model_path,
            "-c", ctx_size,
            "-fa", "on",
            "-ngl", "99",
            "--mmap",
        ]

    def _env_for_role(self, role: str) -> Dict[str, str]:
        if settings.unified_single_model:
            visible = "0,1"   # span both GPUs for the single unified process
        else:
            gpu_idx = self._ROLE_CONFIG[role]["gpu_idx"]
            visible = str(gpu_idx)
        return {
            **os.environ,
            "HIP_VISIBLE_DEVICES":      visible,
            "ROCR_VISIBLE_DEVICES":     visible,
            "HSA_OVERRIDE_GFX_VERSION": "12.0.0",
            "HSA_ENABLE_SDMA":          "0",
            "ROCM_OVERCOMMIT":          "1",
        }

    async def _start_role(self, role: str) -> bool:
        config = self._ROLE_CONFIG[role]
        port = config["port"]
        gpus = "0,1" if settings.unified_single_model else str(config["gpu_idx"])

        async with self._locks[role]:
            await self._terminate_role(role)

            env = self._env_for_role(role)
            cmd = self._build_cmd(role)
            _base_logger.info("Spawning %s (GPU %s / port %d): %s", role, gpus, port, " ".join(cmd))

            try:
                proc = await asyncio.create_subprocess_exec(
                    *cmd,
                    env=env,
                    stdout=asyncio.subprocess.DEVNULL,
                    stderr=asyncio.subprocess.PIPE,
                )
                self._procs[role] = proc
                self._last_activity[role] = time.time()
                # Stream stderr to logger so model errors are visible in proxy logs
                asyncio.create_task(self._pipe_stderr(role, proc))
            except Exception as exc:
                _base_logger.error("Failed to spawn %s: %s", role, exc)
                return False

        if not await self._wait_for_health(config["health_url"], timeout=settings.process_startup_timeout):
            _base_logger.error(
                "%s did not become healthy within %ds", role, settings.process_startup_timeout
            )
            return False

        _base_logger.info("%s ready on port %d (GPU %s)", role, port, gpus)
        return True

    async def _terminate_role(self, role: str) -> None:
        proc = self._procs.get(role)
        if proc is None or proc.returncode is not None:
            self._procs[role] = None
            return
        try:
            proc.terminate()
            await asyncio.wait_for(proc.wait(), timeout=10)
        except asyncio.TimeoutError:
            try:
                proc.kill()
                await proc.wait()
            except ProcessLookupError:
                pass
        except ProcessLookupError:
            pass
        self._procs[role] = None
        _base_logger.info("Terminated %s process", role)

    async def _watchdog(self) -> None:
        """Every 30s: terminate any role idle longer than settings.idle_timeout.
        In unified mode only the single 'reasoner' process is watched."""
        roles_to_watch = ("reasoner",) if settings.unified_single_model else ("reasoner", "coder")
        while True:
            await asyncio.sleep(30)
            now = time.time()
            for role in roles_to_watch:
                idle = now - self._last_activity[role]
                if idle > settings.idle_timeout:
                    proc = self._procs.get(role)
                    if proc is not None and proc.returncode is None:
                        _base_logger.info(
                            "Watchdog: %s idle %.0fs > %ds — terminating to flush VRAM",
                            role, idle, settings.idle_timeout,
                        )
                        async with self._locks[role]:
                            await self._terminate_role(role)

    async def _pipe_stderr(self, role: str, proc: asyncio.subprocess.Process) -> None:
        """Read the process's stderr line-by-line and forward to the proxy logger."""
        stderr_log = logging.getLogger(f"llama-server.{role}")
        assert proc.stderr is not None
        try:
            async for raw in proc.stderr:
                line = raw.decode(errors="replace").rstrip()
                if line:
                    stderr_log.info("%s", line)
        except Exception:
            pass

    async def _wait_for_health(self, health_url: str, timeout: int) -> bool:
        start = time.time()
        async with httpx.AsyncClient(timeout=2) as client:
            while time.time() - start < timeout:
                try:
                    resp = await client.get(health_url)
                    if resp.status_code == 200:
                        return True
                except Exception:
                    pass
                await asyncio.sleep(1.0)
        return False


NATIVE_MODEL_MANAGER = NativeModelManager()

# Persistent circuit breakers — shared across all requests so state is retained.
_REASONER_CIRCUIT = CircuitBreaker(failure_threshold=5, reset_timeout=30.0)
_CODER_CIRCUIT = CircuitBreaker(failure_threshold=5, reset_timeout=30.0)

# ---------------------------------------------------------------------------
# Image job store & cleanup
_image_jobs: Dict[str, Dict[str, Any]] = {}
_background_tasks: Dict[str, asyncio.Task] = {}
_pending_slide_jobs: Dict[str, Dict[str, Any]] = {}   # approval-gated: populated by _finalise
_cleanup_task: Optional[asyncio.Task] = None


async def _cleanup_old_image_jobs() -> None:
    """Periodically remove jobs older than settings.image_job_ttl."""
    while True:
        await asyncio.sleep(300)
        now = time.time()
        to_delete = [
            rid for rid, job in _image_jobs.items()
            if now - job.get("created_at", 0) > settings.image_job_ttl
        ]
        for rid in to_delete:
            _image_jobs.pop(rid, None)
        if to_delete:
            _base_logger.info("Cleaned up %d stale image jobs", len(to_delete))


# ---------------------------------------------------------------------------
# Visual Co-Pilot (batched LLM call for image prompts)
async def enhance_image_prompts(
    slides_data: List[Dict[str, Any]],
    persona: Dict[str, Any],
    client: LLMClient,
    log: RequestLoggerAdapter,
) -> Dict[str, Dict[str, Any]]:
    """One LLM call to enhance image prompts for all slides."""
    image_style = persona.get("image_style", "professional photography")
    palette = persona.get("visual_palette", [])

    batch_input = json_dumps([
        {
            "slide_id": s.get("slide_id", str(i)),
            "title": s.get("title", ""),
            "core_concept": s.get("core_concept") or (s.get("bullets") or [""])[0],
        }
        for i, s in enumerate(slides_data)
    ])

    sys_prompt = (
        "You are a Visual Director specializing in presentation imagery.\n"
        "Given a list of slides, produce ONE enhanced image generation prompt per slide.\n"
        "Style to target: %s\n"
        "Color palette hints: %s\n\n"
        "Rules:\n"
        "- Each prompt must be vivid, descriptive, and under 40 words.\n"
        "- Include composition and lighting directives.\n"
        "- Output ONLY raw valid JSON — an array:\n"
        '  [{"slide_id": "...", "enhanced_prompt": "...", "style_tags": ["tag1", "tag2"]}, ...]\n'
        "No markdown, no fences, no extra keys."
    ) % (image_style, ", ".join(palette))

    try:
        t0 = time.perf_counter()
        data = await client.generate(
            [
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": "Slides:\n" + batch_input},
            ],
            temperature=settings.visual_copilot_temp,
            top_p=settings.visual_copilot_top_p,
            response_format={"type": "json_object"},
        )
        raw = data["choices"][0]["message"]["content"].strip()
        raw = re.sub(r"```(?:json)?\s*\n?", "", raw).strip()
        parsed = json_loads(raw)

        if isinstance(parsed, dict):
            for key in ("slides", "data", "prompts", "items"):
                if key in parsed and isinstance(parsed[key], list):
                    parsed = parsed[key]
                    break
            else:
                for val in parsed.values():
                    if isinstance(val, list):
                        parsed = val
                        break

        if not isinstance(parsed, list):
            raise ValueError("Visual co-pilot response is not a list")

        log.info(
            "Visual co-pilot: %d prompts enhanced in %.2fs",
            len(parsed), time.perf_counter() - t0,
        )
        return {item["slide_id"]: item for item in parsed if "slide_id" in item}
    except Exception as exc:
        log.warning("Visual co-pilot failed, fallback to defaults: %s", exc)
        return {}


# ---------------------------------------------------------------------------
# Background image generation
async def _generate_images_background(
    request_id: str,
    slides: List[Dict[str, Any]],
    enhanced_prompts: Dict[str, Dict[str, Any]],
    log: RequestLoggerAdapter,
) -> None:
    """Generate images for slides that have prompts. Updates _image_jobs."""
    tasks_list: List[Tuple[int, str, str]] = []
    for i, slide in enumerate(slides):
        slide_id = slide.get("slide_id", str(i))
        ep = enhanced_prompts.get(slide_id, {})
        prompt = ep.get("enhanced_prompt") or slide.get("image_prompt", "")
        if prompt:
            tasks_list.append((i, slide_id, prompt))

    job = _image_jobs.setdefault(request_id, {
        "status": "processing",
        "completed": 0,
        "total": len(tasks_list),
        "slides": {},
        "created_at": time.time(),
    })
    job["total"] = len(tasks_list)

    if not tasks_list:
        job["status"] = "done"
        return

    settings.image_output_dir.mkdir(parents=True, exist_ok=True)
    sem = asyncio.Semaphore(settings.image_gen_semaphore)

    async def _gen_one(idx: int, slide_id: str, raw_prompt: str) -> None:
        async with sem:
            filename = settings.image_output_dir / f"slide_{idx}_{uuid.uuid4().hex[:6]}.svg"
            try:
                # SVG generation is CPU-only; run in thread to keep event loop free
                await asyncio.to_thread(generate_svg_pictograph, raw_prompt, filename)
                url = f"/images/{filename.name}"
            except Exception as exc:
                log.warning("SVG generation failed for slide %d: %s", idx, exc)
                url = settings.image_placeholder_url

            job["slides"][slide_id] = url
            job["completed"] += 1
            # Keep the in-memory slide dict up to date so the pending job reflects real URLs
            pending = _pending_slide_jobs.get(request_id, {})
            for slide in pending.get("slides", []):
                if slide.get("slide_id") == slide_id:
                    slide["image_url"] = url
                    break

    await asyncio.gather(*[_gen_one(i, sid, p) for i, sid, p in tasks_list])
    job["status"] = "done"
    log.info("Background image generation complete for request_id=%s", request_id)
    _background_tasks.pop(request_id, None)


# ---------------------------------------------------------------------------
# Reflection loop
async def _reflection_loop(
    final_json: Dict[str, Any],
    schema_str: str,
    persona: Dict[str, Any],
    client: LLMClient,
    params: Dict[str, float],
    log: RequestLoggerAdapter,
) -> Dict[str, Any]:
    """Critic -> revision cycle."""
    for iteration in range(settings.reflection_max_iterations):
        critic_sys = (
            "You are a Presentation Critic. Evaluate the JSON presentation below.\n"
            "Check for:\n"
            "  1. Narrative coherence -- logical story flow\n"
            "  2. Persona adherence -- matches the target profile tone/density\n"
            "  3. Schema compliance -- no critical missing keys\n\n"
            "Persona profile:\n%s\n\n"
            "Output ONLY raw valid JSON (no markdown fences):\n"
            '{"issues": ["..."], "suggestions": ["..."]}\n'
            "If flawless, return empty arrays."
        ) % json_dumps(persona)

        try:
            t0 = time.perf_counter()
            critic_data = await client.generate(
                [
                    {"role": "system", "content": critic_sys},
                    {"role": "user", "content": f"Presentation JSON:\n{json_dumps(final_json)}\n\nSchema:\n{schema_str}"},
                ],
                temperature=0.3,
                top_p=0.5,
                response_format={"type": "json_object"},
            )
            critique = robust_json_loads(critic_data["choices"][0]["message"]["content"])
            issues = critique.get("issues", [])
            suggestions = critique.get("suggestions", [])
            log.info(
                "Critic iteration %d: %d issues, %d suggestions (%.2fs)",
                iteration + 1, len(issues), len(suggestions), time.perf_counter() - t0,
            )
        except Exception as exc:
            log.warning("Critic failed at iteration %d: %s -- skipping.", iteration + 1, exc)
            break

        if not issues:
            log.info("Critic found no issues -- stopping reflection early.")
            break

        revise_sys = (
            "You are a Senior Copywriter revising a presentation based on critic feedback.\n"
            "Apply ALL suggestions. Preserve schema structure exactly.\n"
            "Output ONLY raw valid JSON -- no markdown, no fences."
        )
        try:
            t1 = time.perf_counter()
            rev_data = await client.generate(
                [
                    {"role": "system", "content": revise_sys},
                    {"role": "user", "content": (
                        f"Current JSON:\n{json_dumps(final_json)}\n\n"
                        f"Issues:\n{json_dumps(issues)}\n\n"
                        f"Suggestions:\n{json_dumps(suggestions)}\n\n"
                        f"Schema:\n{schema_str}\n\n"
                        "Return the improved presentation JSON."
                    )},
                ],
                temperature=params["temp"],
                top_p=params["top_p"],
                response_format={"type": "json_object"},
            )
            final_json = robust_json_loads(rev_data["choices"][0]["message"]["content"])
            log.info("Revision %d complete in %.2fs", iteration + 1, time.perf_counter() - t1)
        except Exception as exc:
            log.warning("Revision failed at iteration %d: %s -- keeping previous.", iteration + 1, exc)
            break

    return final_json


# ---------------------------------------------------------------------------
# PPTX compiler
def compile_pptx(presentation_json: dict, output_path: Path) -> None:
    """Compile the pipeline's JSON output into a .pptx file at output_path."""
    prs = pptx.Presentation()

    for slide_data in presentation_json.get("slides", []):
        title_text = slide_data.get("title", "No Title")
        content_items = slide_data.get("content", [])

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text

        text_frame = slide.placeholders[1].text_frame
        text_frame.text = ""

        for i, item in enumerate(content_items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(str(output_path))
    _base_logger.info("PPTX saved: %s", output_path)


# ---------------------------------------------------------------------------
# Orchestrator
class Orchestrator:
    def __init__(self, request_id: str, log: RequestLoggerAdapter) -> None:
        self.request_id = request_id
        self.log = log
        self.reasoner_client = LLMClient(settings.reasoner_url, timeout=settings.llm_timeout, circuit=_REASONER_CIRCUIT)
        # In unified mode both phases share the same endpoint and circuit breaker
        coder_url = settings.reasoner_url if settings.unified_single_model else settings.coder_url
        coder_circuit = _REASONER_CIRCUIT if settings.unified_single_model else _CODER_CIRCUIT
        self.coder_client = LLMClient(coder_url, timeout=settings.llm_timeout, circuit=coder_circuit)

    async def run_full(
        self,
        chat_req: ChatCompletionRequest,
        params: Dict[str, float],
        persona_key: str,
        narrative: NarrativeStructure,
        pipeline_mode: PipelineMode,
    ) -> Dict[str, Any]:
        """Full pipeline: Supervisor (reasoner / GPU 0) -> Worker (coder / GPU 1)."""
        schema, schema_str = self._extract_schema(chat_req.response_format)
        persona = PERSONA_PROFILES.get(persona_key, PERSONA_PROFILES["balanced"])

        # Phase 1: Reasoner — ensure process is alive before calling
        self.log.info("Phase 1: routing to reasoner (GPU 0 / port 8081)")
        if not await NATIVE_MODEL_MANAGER.ensure_role_ready("reasoner"):
            raise_api_error(ErrorCode.SUPERVISOR_FAILED, "Reasoner process could not start.", 502)
        outline, arc = await self._run_supervisor(chat_req, params, persona, narrative, schema_str)

        if pipeline_mode == PipelineMode.OUTLINE_ONLY:
            await self.close()
            return {
                "id": f"presenton-{self.request_id}",
                "object": "chat.completion",
                "presenton": {"outline": outline, "presentation_arc": arc},
            }

        # Phase 2: Coder — ensure process is alive before calling
        self.log.info("Phase 2: routing to coder (GPU 1 / port 8082)")
        if not await NATIVE_MODEL_MANAGER.ensure_role_ready("coder"):
            raise_api_error(ErrorCode.WORKER_FAILED, "Coder process could not start.", 502)
        final_json, worker_data = await self._run_worker(outline, arc, schema_str, persona, params)
        result = await self._finalise(final_json, worker_data, schema, schema_str, persona, persona_key, narrative, params)
        await self.close()
        return result

    async def run_expand(self, expand_req: ExpandRequest, params: Dict[str, float]) -> Dict[str, Any]:
        """Expand a user-supplied outline directly using the coder (GPU 1)."""
        persona_key = expand_req.persona.lower()
        persona = PERSONA_PROFILES.get(persona_key, PERSONA_PROFILES["balanced"])
        schema = expand_req.response_format.get("json_schema", {}).get("schema", {})
        schema_str = json_dumps(schema) if schema else "{}"
        narrative = parse_enum_ci(
            NarrativeStructure,
            expand_req.narrative_structure,
            NarrativeStructure.PROBLEM_SOLUTION,
            log_warning=True,
        )
        arc = {"note": "User-supplied outline."}

        self.log.info("Direct expansion: routing to coder (GPU 1 / port 8082)")
        if not await NATIVE_MODEL_MANAGER.ensure_role_ready("coder"):
            raise_api_error(ErrorCode.WORKER_FAILED, "Coder process could not start.", 502)
        final_json, worker_data = await self._run_worker(expand_req.outline, arc, schema_str, persona, params)
        result = await self._finalise(final_json, worker_data, schema, schema_str, persona, persona_key, narrative, params)
        await self.close()
        return result

    def _extract_schema(self, response_format: Dict[str, Any]) -> Tuple[Dict, str]:
        if response_format.get("type") != "json_schema":
            raise_api_error(ErrorCode.UNSUPPORTED_FORMAT, "response_format.type must be 'json_schema'.", 400)
        schema = response_format.get("json_schema", {}).get("schema")
        if not schema:
            raise_api_error(ErrorCode.INVALID_REQUEST, "Missing schema in response_format.", 400)
        schema_str = json_dumps(schema)
        if len(schema_str.encode()) > settings.max_schema_size:
            raise_api_error(ErrorCode.REQUEST_TOO_LARGE, "Schema payload too large.", 413)
        return schema, schema_str

    async def _keepalive_touch(self, role: str, stop_event: asyncio.Event) -> None:
        """Pings the watchdog every 10 s while a request is in flight."""
        while not stop_event.is_set():
            NATIVE_MODEL_MANAGER.touch(role)
            await asyncio.sleep(10)

    async def _run_supervisor(
        self,
        chat_req: ChatCompletionRequest,
        params: Dict[str, float],
        persona: Dict[str, Any],
        narrative: NarrativeStructure,
        schema_str: str,
    ) -> Tuple[List[Dict], Dict]:
        narrative_instr = NARRATIVE_INSTRUCTIONS.get(narrative, "")
        persona_summary = (
            f"Persona: {persona.get('tone', '')} | "
            f"Tone: {persona.get('sentence_style', '')} | "
            f"Density: {persona.get('slide_density', '')}"
        )

        supervisor_sys = (
            "You are the Lead Presentation Architect.\n\n"
            f"## Persona Mode\n{persona_summary}\n\n"
            f"## Narrative Arc Rules\n{narrative_instr}\n\n"
            "## Requirements\n"
            "Output an architectural presentation outline JSON exactly matching this shape:\n"
            "{\n"
            '  "presentation_arc": "...",\n'
            '  "emotional_curve": "...",\n'
            '  "outline": [\n'
            '    {\n'
            '      "slide_id": "<uuid4>",\n'
            '      "title": "...",\n'
            '      "core_concept": "...",\n'
            '      "bullet_points": ["..."],\n'
            '      "story_beat": "..."\n'
            '    }\n'
            "  ]\n"
            "}\n\n"
            "Rules:\n"
            "- Assign a unique UUID4 to slide_id.\n"
            "- Do NOT generate final prose.\n"
            "- Output ONLY raw valid JSON -- no markdown fences.\n\n"
            f"Target Output Schema:\n{schema_str}"
        )

        msgs = [{"role": "system", "content": supervisor_sys}]
        for msg in chat_req.messages[-settings.max_messages:]:
            if msg.role == "system":
                continue
            content = sanitize_prompt(msg.content) if msg.role == "user" else msg.content
            msgs.append({"role": msg.role, "content": content})

        t0 = time.perf_counter()
        stop_event = asyncio.Event()
        heartbeat_task = asyncio.create_task(self._keepalive_touch("reasoner", stop_event))
        try:
            sup_data = await self.reasoner_client.generate(
                msgs,
                temperature=params["temp"],
                top_p=params["top_p"],
                response_format={"type": "json_object"},
            )
            parsed = robust_json_loads(sup_data["choices"][0]["message"]["content"])
            outline = parsed.get("outline", [])
            arc = {
                "presentation_arc": parsed.get("presentation_arc", ""),
                "emotional_curve": parsed.get("emotional_curve", ""),
            }
            self.log.info(
                "Supervisor: %d slides, arc='%s' (%.2fs)",
                len(outline), arc.get("presentation_arc", "")[:60], time.perf_counter() - t0,
            )
            return outline, arc
        except RetryError:
            raise_api_error(ErrorCode.SUPERVISOR_FAILED, "Supervisor unavailable after retries.", 502)
        except HTTPException:
            raise
        except Exception as exc:
            self.log.exception("Supervisor failure")
            raise_api_error(ErrorCode.SUPERVISOR_FAILED, f"Supervisor failed: {exc}", 502)
        finally:
            stop_event.set()
            heartbeat_task.cancel()
            try:
                await heartbeat_task
            except asyncio.CancelledError:
                pass

    async def _run_worker(
        self,
        outline: List[Dict],
        arc: Dict,
        schema_str: str,
        persona: Dict[str, Any],
        params: Dict[str, float],
    ) -> Tuple[Dict, Dict]:
        persona_detail = json_dumps({
            k: persona[k]
            for k in ("tone", "sentence_style", "slide_density", "layout_bias")
            if k in persona
        })
        layout_registry = json_dumps(LAYOUT_DEFINITIONS)
        typography = json_dumps(persona.get("typography", {"title": "Inter", "body": "Source Sans Pro"}))
        palette = json_dumps(persona.get("visual_palette", []))

        worker_sys = (
            "You are a Creative Copywriter and Art Director.\n\n"
            f"## Persona Specifications\n{persona_detail}\n\n"
            f"## Structure Arc Parameters\n{json_dumps(arc)}\n\n"
            f"## Layout Key Library\n{layout_registry}\n\n"
            "## Tasks\n"
            "1. Expand each outline node into engaging, on-persona prose.\n"
            "2. Preserve schema structure exactly.\n"
            "3. For each slide, populate `_visual_hints` as:\n"
            '   {"icon": "...", "layout": {"type": "...", "image_weight": ..., "text_weight": ...}, '
            '"image_prompt": "...", "theme": "default", '
            f'"font_pair": {typography}, "colors": {palette}}}\n'
            "4. Carry slide_id and story_beat through unchanged.\n"
            "5. Output ONLY raw valid JSON -- no markdown, no fences."
        )

        t0 = time.perf_counter()
        stop_event = asyncio.Event()
        heartbeat_task = asyncio.create_task(self._keepalive_touch("coder", stop_event))
        try:
            worker_data = await self.coder_client.generate(
                [
                    {"role": "system", "content": worker_sys},
                    {"role": "user", "content": (
                        f"Outline:\n{json_dumps(outline)}\n\n"
                        f"Target Schema:\n{schema_str}\n\n"
                        "Produce the final presentation JSON."
                    )},
                ],
                temperature=params["temp"],
                top_p=params["top_p"],
                response_format={"type": "json_object"},
            )
            final_json = robust_json_loads(worker_data["choices"][0]["message"]["content"])
            self.log.info("Worker completed in %.2fs", time.perf_counter() - t0)
            return final_json, worker_data
        except RetryError:
            raise_api_error(ErrorCode.WORKER_FAILED, "Worker unavailable after retries.", 502)
        except HTTPException:
            raise
        except Exception as exc:
            self.log.exception("Worker failure")
            raise_api_error(ErrorCode.WORKER_FAILED, f"Worker failed: {exc}", 502)
        finally:
            stop_event.set()
            heartbeat_task.cancel()
            try:
                await heartbeat_task
            except asyncio.CancelledError:
                pass

    async def _finalise(
        self,
        final_json: Dict[str, Any],
        worker_data: Dict[str, Any],
        schema: Dict,
        schema_str: str,
        persona: Dict[str, Any],
        persona_key: str,
        narrative: NarrativeStructure,
        params: Dict[str, float],
    ) -> Dict[str, Any]:
        # Reflection (if enabled and persona allows)
        if settings.reflection_enabled and persona.get("reflection_gate", False):
            final_json = await _reflection_loop(
                final_json, schema_str, persona, self.coder_client, params, self.log
            )

        # Schema validation
        if schema:
            try:
                validate_against_schema(final_json, schema)
                self.log.info("Output validated against schema.")
            except ValueError as exc:
                self.log.error("Schema validation error: %s", exc)
                raise_api_error(ErrorCode.SCHEMA_VALIDATION_FAILED, str(exc), 422)

        # Find slides list (robust depth-1 search)
        slides_list = self._find_slides_list(final_json)

        enhanced_prompts: Dict[str, Dict] = {}
        if settings.image_gen_enabled and slides_list:
            # Enhance prompts via visual co-pilot so SVGs are semantically informed
            enhanced_prompts = await enhance_image_prompts(slides_list, persona, self.coder_client, self.log)
            for slide in slides_list:
                sid = slide.get("slide_id")
                if sid and sid in enhanced_prompts:
                    ep = enhanced_prompts[sid]
                    hints = slide.setdefault("_visual_hints", {})
                    hints["enhanced_prompt"] = ep.get("enhanced_prompt")
                    hints["style_tags"] = ep.get("style_tags", [])
                # Images are null until the user approves via /accept
                slide["image_url"] = None

            # Queue for approval-gated generation; no background task fires yet
            _pending_slide_jobs[self.request_id] = {
                "status": "pending_approval",
                "slides": slides_list,
                "enhanced_prompts": enhanced_prompts,
                "persona": persona_key,
                "narrative": narrative.value,
                "created_at": time.time(),
            }

        # Compile PPTX in a thread so blocking file I/O doesn't stall the event loop
        settings.pptx_output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = settings.pptx_output_dir / f"presentation_{self.request_id}.pptx"
        await asyncio.to_thread(compile_pptx, final_json, pptx_path)

        # Metadata
        images_pending = (
            settings.image_gen_enabled
            and bool(slides_list)
            and self.request_id in _pending_slide_jobs
        )
        final_json["_meta"] = {
            "pipeline_version": "6.0.0 (Native Sequential)",
            "supervisor_model": f"Reasoner ({settings.reasoner_model_key}) on GPU 0",
            "worker_model": f"Coder ({settings.coder_model_key}) on GPU 1",
            "persona": persona_key,
            "narrative_structure": narrative.value,
            "timestamp": time.time(),
            "request_id": self.request_id,
            "download_url": f"/v1/presenton/download/{self.request_id}",
            "image_status": "pending_approval" if images_pending else "disabled",
            "approval_url": f"/v1/presenton/accept/{self.request_id}" if images_pending else None,
        }

        worker_data["choices"][0]["message"]["content"] = json_dumps(final_json)
        return worker_data

    def _find_slides_list(self, data: Dict[str, Any]) -> List[Dict]:
        """Search for slides list top-level, then one level deeper, skipping metadata keys."""
        for key in settings.slides_list_keys:
            val = data.get(key)
            if isinstance(val, list):
                return val

        for k, v in data.items():
            if k.startswith("_") or k in settings.skip_keys_depth1:
                continue
            if isinstance(v, dict):
                for key in settings.slides_list_keys:
                    sub = v.get(key)
                    if isinstance(sub, list):
                        return sub
        return []

    async def close(self) -> None:
        await self.reasoner_client.close()
        await self.coder_client.close()


# ---------------------------------------------------------------------------
# Lifespan (startup & shutdown)
@asynccontextmanager
async def lifespan(app: FastAPI):
    global _cleanup_task
    _base_logger.info("Presenton proxy starting up (Native Sequential mode)...")

    ready = await NATIVE_MODEL_MANAGER.start_all()
    if not ready:
        _base_logger.warning("One or both model processes failed to become ready – check logs.")
    else:
        _base_logger.info("Both model processes are ready and healthy.")

    _cleanup_task = asyncio.create_task(_cleanup_old_image_jobs())

    if settings.image_gen_enabled:
        settings.image_output_dir.mkdir(parents=True, exist_ok=True)
        app.mount("/images", StaticFiles(directory=str(settings.image_output_dir)), name="images")

    yield

    _base_logger.info("Shutting down proxy...")
    if _cleanup_task:
        _cleanup_task.cancel()
        try:
            await _cleanup_task
        except asyncio.CancelledError:
            pass

    if _background_tasks:
        _base_logger.info("Waiting for %d pending image generation tasks...", len(_background_tasks))
        _, pending = await asyncio.wait(_background_tasks.values(), timeout=30)
        if pending:
            _base_logger.warning("%d image generation tasks did not finish within 30s", len(pending))

    await NATIVE_MODEL_MANAGER.shutdown()
    _base_logger.info("Shutdown complete.")


# ---------------------------------------------------------------------------
# FastAPI application
app = FastAPI(
    title="Presenton Native Multi-Agent Proxy",
    version="6.0.0",
    lifespan=lifespan,
)

# CORS middleware — localhost:3000 always trusted; add more via CORS_ORIGINS env var
_cors_origins = list({
    "http://localhost:3000",
    *[o.strip() for o in settings.cors_origins.split(",") if o.strip()],
})
app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.exception_handler(HTTPException)
async def _global_http_exc_handler(request: Request, exc: HTTPException) -> JSONResponse:
    request_id = request.headers.get("x-request-id", str(uuid.uuid4()))
    return _http_exc_to_response(exc, request_id)


def _resolve_params(mode_str: Optional[str]) -> Tuple[HallucinationMode, Dict[str, float]]:
    mode = parse_enum_ci(
        HallucinationMode,
        mode_str or settings.hallucination_mode,
        HallucinationMode.TIGHT,
        log_warning=True,
    )
    if mode == HallucinationMode.ACCOMMODATIVE:
        return mode, {"temp": settings.temp_accommodative, "top_p": settings.top_p_accommodative}
    return mode, {"temp": settings.temp_tight, "top_p": settings.top_p_tight}


# ---------------------------------------------------------------------------
# API endpoints
@app.post("/v1/chat/completions")
async def proxy_chat_completions(
    request: Request,
    x_hallucination_mode: Optional[str] = Header(None, alias="x-hallucination-mode"),
    x_pipeline_mode: Optional[str] = Header(None, alias="x-pipeline-mode"),
    x_persona: Optional[str] = Header(None, alias="x-persona"),
    x_narrative_structure: Optional[str] = Header(None, alias="x-narrative-structure"),
    x_request_id: Optional[str] = Header(None, alias="x-request-id"),
    _auth: None = Depends(require_api_key),
) -> JSONResponse:
    request_id = x_request_id or str(uuid.uuid4())
    log = make_logger(request_id)
    t_start = time.perf_counter()
    error_occurred = False

    try:
        raw_body = await request.body()
        if len(raw_body) > settings.max_body_bytes:
            return error_response(413, ErrorCode.REQUEST_TOO_LARGE, "Request body too large.", request_id)

        try:
            chat_req = ChatCompletionRequest.model_validate(json_loads(raw_body))
        except (ValidationError, ValueError, json.JSONDecodeError) as exc:
            return error_response(422, ErrorCode.INVALID_REQUEST, str(exc), request_id)

        _, params = _resolve_params(x_hallucination_mode)
        pipeline_mode = parse_enum_ci(PipelineMode, x_pipeline_mode or "full", PipelineMode.FULL, log_warning=True)
        persona_key = (x_persona or settings.default_persona).lower()
        narrative = parse_enum_ci(
            NarrativeStructure,
            x_narrative_structure or settings.default_narrative,
            NarrativeStructure.PROBLEM_SOLUTION,
            log_warning=True,
        )

        log.info(
            "Request: pipeline=%s persona=%s narrative=%s temp=%.2f",
            pipeline_mode.value, persona_key, narrative.value, params["temp"],
        )

        orch = Orchestrator(request_id, log)
        result = await orch.run_full(chat_req, params, persona_key, narrative, pipeline_mode)
        return JSONResponse(content=result, media_type="application/json", headers={"X-Request-ID": request_id})

    except HTTPException as exc:
        error_occurred = True
        return _http_exc_to_response(exc, request_id)
    except Exception:
        error_occurred = True
        log.exception("Unhandled pipeline error")
        return error_response(500, ErrorCode.INTERNAL_ERROR, "Internal server error.", request_id)
    finally:
        latency = time.perf_counter() - t_start
        _REQUEST_LATENCY.observe(latency)
        _REQUESTS_TOTAL.labels(endpoint="chat_completions", status="error" if error_occurred else "ok").inc()
        log.info("Completed in %.3fs (error=%s)", latency, error_occurred)


@app.post("/v1/presenton/expand")
async def expand_outline(
    request: Request,
    x_hallucination_mode: Optional[str] = Header(None, alias="x-hallucination-mode"),
    x_request_id: Optional[str] = Header(None, alias="x-request-id"),
    _auth: None = Depends(require_api_key),
) -> JSONResponse:
    request_id = x_request_id or str(uuid.uuid4())
    log = make_logger(request_id)
    t_start = time.perf_counter()
    error_occurred = False

    try:
        raw_body = await request.body()
        if len(raw_body) > settings.max_body_bytes:
            return error_response(413, ErrorCode.REQUEST_TOO_LARGE, "Request body too large.", request_id)

        try:
            expand_req = ExpandRequest.model_validate(json_loads(raw_body))
        except (ValidationError, ValueError, json.JSONDecodeError) as exc:
            return error_response(422, ErrorCode.INVALID_REQUEST, str(exc), request_id)

        _, params = _resolve_params(x_hallucination_mode)
        log.info("Expand: persona=%s slides=%d", expand_req.persona, len(expand_req.outline))

        orch = Orchestrator(request_id, log)
        result = await orch.run_expand(expand_req, params)
        return JSONResponse(content=result, media_type="application/json", headers={"X-Request-ID": request_id})

    except HTTPException as exc:
        error_occurred = True
        return _http_exc_to_response(exc, request_id)
    except Exception:
        error_occurred = True
        log.exception("Unhandled expand error")
        return error_response(500, ErrorCode.INTERNAL_ERROR, "Internal server error.", request_id)
    finally:
        latency = time.perf_counter() - t_start
        _REQUEST_LATENCY.observe(latency)
        _REQUESTS_TOTAL.labels(endpoint="expand", status="error" if error_occurred else "ok").inc()
        log.info("Expand completed in %.3fs", latency)


@app.get("/v1/presenton/images/status")
async def image_status(
    request_id: str = Query(..., description="Request ID from original completion"),
) -> JSONResponse:
    job = _image_jobs.get(request_id)
    if job is None:
        # Image generation disabled or not yet started — report ready with zero images
        return JSONResponse(content={
            "ready": True,
            "completed": 0,
            "total": 0,
            "slides": {},
        })
    return JSONResponse(content={
        "ready": job["status"] == "done",
        "completed": job["completed"],
        "total": job["total"],
        "slides": job["slides"],
    })


@app.post("/v1/presenton/accept/{request_id}")
async def accept_presentation(request_id: str) -> JSONResponse:
    job = _pending_slide_jobs.get(request_id)
    if job is None:
        raise HTTPException(
            status_code=404,
            detail={"code": ErrorCode.NOT_FOUND.value, "message": f"No pending job for request_id '{request_id}'."},
        )
    if job["status"] != "pending_approval":
        return JSONResponse(content={
            "request_id": request_id,
            "status": job["status"],
            "message": "Job is not awaiting approval.",
        })
    job["status"] = "generating"
    log = make_logger(request_id)
    task = asyncio.create_task(
        _generate_images_background(request_id, job["slides"], job["enhanced_prompts"], log)
    )
    _background_tasks[request_id] = task
    return JSONResponse(content={
        "request_id": request_id,
        "status": "generating",
        "message": "SVG image generation started. Poll /v1/presenton/images/status for progress.",
    })


@app.get("/v1/presenton/download/{request_id}")
async def download_pptx(request_id: str) -> FileResponse:
    pptx_path = settings.pptx_output_dir / f"presentation_{request_id}.pptx"
    if not pptx_path.exists():
        raise HTTPException(
            status_code=404,
            detail={"code": ErrorCode.NOT_FOUND.value, "message": f"No PPTX found for request_id '{request_id}'."},
        )
    return FileResponse(
        path=str(pptx_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"presentation_{request_id}.pptx",
    )


# ---------------------------------------------------------------------------
# Health & observability
@app.get("/health")
async def health() -> Dict:
    return {"status": "ok", "version": "6.0.0", "engine_mode": "native_sequential_dual_gpu"}


@app.get("/health/detailed")
async def health_detailed() -> Dict:
    async def check_url(url: str) -> bool:
        try:
            async with httpx.AsyncClient(timeout=3) as c:
                r = await c.get(url)
                return r.status_code == 200
        except Exception:
            return False

    reasoner_ok = await check_url(settings.reasoner_health_url)
    coder_ok = await check_url(settings.coder_health_url)
    proc_status = NATIVE_MODEL_MANAGER.process_status()

    return {
        "status": "ok" if (reasoner_ok and coder_ok) else "degraded",
        "proxy": "ok",
        "reasoner": "ok" if reasoner_ok else "unreachable",
        "coder": "ok" if coder_ok else "unreachable",
        "processes": proc_status,
        "circuit_breakers": {
            "reasoner": _REASONER_CIRCUIT.state,
            "coder": _CODER_CIRCUIT.state,
        },
    }


@app.get("/metrics")
async def get_metrics() -> Response:
    return Response(generate_latest(), media_type=CONTENT_TYPE_LATEST)


@app.get("/v1/presenton/personas")
async def list_personas() -> Dict:
    public_fields = (
        "tone", "sentence_style", "slide_density", "layout_bias",
        "rhetorical_devices", "typography", "visual_palette", "image_style",
    )
    return {
        key: {f: profile[f] for f in public_fields if f in profile}
        for key, profile in PERSONA_PROFILES.items()
    }


@app.get("/v1/presenton/layouts")
async def list_layouts() -> Dict:
    return LAYOUT_DEFINITIONS


@app.get("/v1/presenton/narratives")
async def list_narratives() -> Dict:
    return {ns.value: NARRATIVE_INSTRUCTIONS[ns] for ns in NarrativeStructure}


# ---------------------------------------------------------------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("presenton_proxy:app", host="0.0.0.0", port=8000, log_level="info")
